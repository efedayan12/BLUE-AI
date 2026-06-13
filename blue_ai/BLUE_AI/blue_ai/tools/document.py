"""
BLUE_AI — Belge Olusturma Araci

Word, Excel, PowerPoint belgeleri olusturur.
LLM ile icerik uretir, python-docx/openpyxl/python-pptx ile dosyaya yazar.
"""

import os
from pathlib import Path
from datetime import datetime

from blue_ai.ai.llm import get_llm


DOCS_DIR = Path.home() / "Documents" / "BLUE_AI"
DOCS_DIR.mkdir(parents=True, exist_ok=True)


def create_word_document(topic: str) -> dict:
    """Word belgesi olustur."""
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {"success": False, "error": "python-docx yuklu degil. 'pip install python-docx' calistirin."}

    llm = get_llm()

    # LLM'den icerik iste (yoksa template fallback)
    if llm.is_available():
        system = """Sen bir profesyonel belge yazarisin. Kullanicinin istedigi konuda 
yapilandirilmis bir belge icerigi olustur. Baslik, alt basliklar ve paragraflar halinde yaz.
Her bolumu '## Baslik' formatinda yaz. Turkce yaz."""
        content = llm.generate(f"Su konuda detayli bir belge hazirla: {topic}", system=system)
    else:
        # Template fallback — LLM olmadan temel yapi
        content = f"""## Giris
{topic} hakkinda hazirlanan bu belge, konunun temel noktalarini ele almaktadir.

## Kapsam
Bu bolumde {topic} konusunun kapsami aciklanmaktadir.

## Detaylar
{topic} ile ilgili detayli bilgiler bu bolumde yer almaktadir.

## Sonuc
{topic} konusu ozetlenmis ve degerlendirme yapilmistir.

## Notlar
- Belge BLUE_AI tarafindan otomatik olusturulmustur
- Icerik duzenlemesi yapilabilir
- Tarih: {datetime.now().strftime('%d.%m.%Y')}"""

    # Word belgesi olustur
    doc = Document()

    # Baslik
    title = doc.add_heading(level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(topic.title())
    run.font.size = Pt(24)

    # Tarih
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_run = date_para.add_run(f"Olusturma Tarihi: {datetime.now().strftime('%d.%m.%Y')}")
    date_run.font.size = Pt(10)

    doc.add_paragraph()  # Bosluk

    # Icerigi isle
    lines = content.split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- ") or line.startswith("* "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)

    # Kaydet
    safe_name = "".join(c for c in topic if c.isalnum() or c in " _-").strip()[:50]
    filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    filepath = DOCS_DIR / filename
    doc.save(str(filepath))

    return {
        "success": True,
        "path": str(filepath),
        "filename": filename,
        "content_length": len(content),
    }


def create_excel_spreadsheet(topic: str) -> dict:
    """Excel tablosu olustur."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    except ImportError:
        return {"success": False, "error": "openpyxl yuklu degil. 'pip install openpyxl' calistirin."}

    llm = get_llm()

    if llm.is_available():
        system = """Sen bir veri analisti ve tablo uzmansin. Kullanicinin istedigi konuda
bir Excel tablosu icin veri olustur. Ciktiyi su formatta ver:
BASLIK: Tablo basligi
SUTUNLAR: sutun1 | sutun2 | sutun3
VERI: deger1 | deger2 | deger3
VERI: deger4 | deger5 | deger6
(en az 5-10 satir veri olustur)"""
        content = llm.generate(f"Su konu icin detayli bir tablo hazirla: {topic}", system=system)
    else:
        content = f"""BASLIK: {topic}
SUTUNLAR: No | Bilgi | Deger | Notlar
VERI: 1 | Konu | {topic} | BLUE_AI
VERI: 2 | Tarih | {datetime.now().strftime('%d.%m.%Y')} | Otomatik
VERI: 3 | Durum | Aktif | -
VERI: 4 | Oncelik | Yuksek | -
VERI: 5 | Aciklama | Detay eklenecek | -"""

    wb = Workbook()
    ws = wb.active
    ws.title = topic[:30]

    header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = Font(name="Calibri", size=12, bold=True, color="FFFFFF")
    cell_font = Font(name="Calibri", size=11)
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    row_idx = 1
    lines = content.split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if line.startswith("BASLIK:"):
            ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
            cell = ws.cell(row=1, column=1, value=line[7:].strip())
            cell.font = Font(name="Calibri", size=16, bold=True)
            cell.alignment = Alignment(horizontal="center")
            row_idx = 3

        elif line.startswith("SUTUNLAR:"):
            cols = [c.strip() for c in line[9:].split("|")]
            for ci, col in enumerate(cols, 1):
                cell = ws.cell(row=row_idx, column=ci, value=col)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border
                ws.column_dimensions[chr(64 + ci)].width = 20
            row_idx += 1

        elif line.startswith("VERI:"):
            vals = [v.strip() for v in line[5:].split("|")]
            for ci, val in enumerate(vals, 1):
                cell = ws.cell(row=row_idx, column=ci, value=val)
                cell.font = cell_font
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="center")
            row_idx += 1

    safe_name = "".join(c for c in topic if c.isalnum() or c in " _-").strip()[:50]
    filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    filepath = DOCS_DIR / filename
    wb.save(str(filepath))

    return {"success": True, "path": str(filepath), "filename": filename}


def create_presentation(topic: str) -> dict:
    """PowerPoint sunumu olustur."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"success": False, "error": "python-pptx yuklu degil. 'pip install python-pptx' calistirin."}

    llm = get_llm()

    if llm.is_available():
        system = """Sen bir sunum uzmansin. Kullanicinin istedigi konuda PowerPoint sunumu icin
icerik olustur. Her slayt icin su formati kullan:
SLAYT: Slayt Basligi
ICERIK: Madde 1
ICERIK: Madde 2
ICERIK: Madde 3
(En az 5-7 slayt olustur)"""
        content = llm.generate(f"Su konu icin sunum hazirla: {topic}", system=system)
    else:
        content = f"""SLAYT: Giris
ICERIK: {topic} hakkinda genel bilgi
ICERIK: Konunun onemi ve kapsami
ICERIK: BLUE_AI tarafindan hazirlanmistir
SLAYT: Kapsam
ICERIK: Temel noktalar
ICERIK: Hedefler ve beklentiler
ICERIK: Zaman plani
SLAYT: Detaylar
ICERIK: Teknik bilgiler
ICERIK: Uygulama adimlari
ICERIK: Kaynaklar
SLAYT: Sonuc
ICERIK: Ozet ve degerlendirme
ICERIK: Sonraki adimlar
ICERIK: Sorular"""

    prs = Presentation()

    # Kapak slayt
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic.title()
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"BLUE_AI | {datetime.now().strftime('%d.%m.%Y')}"

    # Icerik slaylar
    current_title = ""
    current_items = []

    lines = content.split("\n")
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if line.startswith("SLAYT:"):
            if current_title and current_items:
                _add_slide(prs, current_title, current_items)
            current_title = line[6:].strip()
            current_items = []
        elif line.startswith("ICERIK:"):
            current_items.append(line[7:].strip())

    if current_title and current_items:
        _add_slide(prs, current_title, current_items)

    safe_name = "".join(c for c in topic if c.isalnum() or c in " _-").strip()[:50]
    filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    filepath = DOCS_DIR / filename
    prs.save(str(filepath))

    return {"success": True, "path": str(filepath), "filename": filename}


def _add_slide(prs, title: str, items: list[str]):
    """Sunum slayt ekle."""
    from pptx.util import Pt
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            tf.text = item
        else:
            p = tf.add_paragraph()
            p.text = item
        tf.paragraphs[-1].font.size = Pt(18)
